"""
Build Account Research - Quick Start.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ── Constants ────────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

C_DARK_PURPLE   = RGBColor(0x26, 0x06, 0x5D)
C_PRIMARY       = RGBColor(0x4C, 0x00, 0xFF)
C_LIGHT_PURPLE  = RGBColor(0xEE, 0xEA, 0xFF)
C_RED           = RGBColor(0xFF, 0x52, 0x52)
C_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT     = RGBColor(0x11, 0x18, 0x27)
C_MUTED_TEXT    = RGBColor(0x6B, 0x72, 0x80)
C_CREAM_BG      = RGBColor(0xFF, 0xFB, 0xEB)
C_CREAM_BORDER  = RGBColor(0xFD, 0xE6, 0x8A)
C_AMBER         = RGBColor(0xD9, 0x77, 0x06)
C_ROW_ALT       = RGBColor(0xF8, 0xF3, 0xF0)
C_BORDER_LIGHT  = RGBColor(0xE0, 0xD9, 0xFF)
C_LIGHT_GRAY    = RGBColor(0xE5, 0xE7, 0xEB)
C_F9            = RGBColor(0xF9, 0xF9, 0xF9)

SLIDE_W_IN = 13.33
SLIDE_H_IN = 7.5

EMU = 914400  # per inch


def emu(inches):
    return int(inches * EMU)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        emu(left), emu(top), emu(width), emu(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None,
                     line_color=None, line_width=1.0, radius_pct=5):
    """Add a rounded rectangle."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE
        emu(left), emu(top), emu(width), emu(height)
    )
    # Set corner radius via XML adj attribute
    sp = shape.element
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        # Clear existing
        for gd in avLst.findall(qn('a:gd')):
            avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius_pct * 333}')

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_name='Calibri',
                font_size=14, bold=False, italic=False, color=None,
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a simple single-run textbox."""
    txb = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def set_shape_text_multiline(shape, lines, font_name='Calibri', font_size=14,
                              bold=False, italic=False, color=None,
                              align=PP_ALIGN.LEFT, para_space_after=Pt(4)):
    """Set multi-line text on an existing shape (clears existing paragraphs)."""
    tf = shape.text_frame
    tf.word_wrap = True
    # Clear
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    for idx, line in enumerate(lines):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        if para_space_after:
            para.space_after = para_space_after
        run = para.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color


def add_multiline_textbox(slide, left, top, width, height, lines,
                           font_name='Calibri', font_size=14,
                           bold=False, italic=False, color=None,
                           align=PP_ALIGN.LEFT):
    """Add a textbox with multiple paragraphs."""
    txb = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
    return txb


def add_header_bars(slide):
    """Add the dark purple + thin red top bar used on content slides."""
    add_rect(slide, 0, 0, SLIDE_W_IN, 0.09, fill_color=C_DARK_PURPLE)
    add_rect(slide, 0, 0.09, SLIDE_W_IN, 0.045, fill_color=C_RED)


def set_slide_bg_white(slide, prs):
    """Ensure slide background is explicitly white."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_WHITE


def set_slide_bg_dark(slide, prs):
    """Ensure slide background is dark purple."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_DARK_PURPLE


def add_table_slide(slide, left, top, width, col_widths, rows_data,
                    header_fill=None, alt_fill=None, default_fill=None,
                    font_size=12):
    """
    Add a table to the slide.
    rows_data: list of tuples — (cells_list, bold_first_col, is_header)
    col_widths: list of widths in inches
    """
    ncols = len(col_widths)
    nrows = len(rows_data)
    col_w_emu = [emu(w) for w in col_widths]
    total_w = sum(col_widths)
    row_height = emu(0.42)

    tbl = slide.shapes.add_table(nrows, ncols, emu(left), emu(top),
                                  emu(total_w), row_height * nrows)
    table = tbl.table

    # Set column widths
    for ci, cw in enumerate(col_w_emu):
        table.columns[ci].width = cw

    for ri, (cells, bold_first, is_header) in enumerate(rows_data):
        row = table.rows[ri]
        row.height = row_height
        for ci, cell_text in enumerate(cells):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run() if not para.runs else para.runs[0]
            run.text = cell_text
            run.font.name = 'Calibri'
            run.font.size = Pt(font_size)
            is_bold = bold_first and ci == 0
            run.font.bold = is_header or is_bold

            if is_header:
                run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill or C_DARK_PURPLE
            else:
                run.font.color.rgb = C_DARK_TEXT
                if ri % 2 == 0 and alt_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = default_fill or C_WHITE

            # Set cell border (light gray)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = etree.SubElement(tcPr, qn(border_tag))
                ln.set('w', '6350')  # 0.5pt
                ln.set('cap', 'flat')
                ln.set('cmpd', 'sng')
                solidFill = etree.SubElement(ln, qn('a:solidFill'))
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
                srgbClr.set('val', 'D1D5DB')

    return tbl


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_cover(slide, prs):
    set_slide_bg_dark(slide, prs)

    # Logo text
    add_textbox(slide, 0.45, 0.35, 3.0, 0.4,
                'D  docusign', font_size=18, bold=True, color=C_WHITE)

    # Genius Bar badge (pill)
    badge = add_rounded_rect(slide, 3.2, 0.35, 1.4, 0.35,
                              fill_color=C_LIGHT_PURPLE, line_color=None, radius_pct=15)
    txb = slide.shapes.add_textbox(emu(3.2), emu(0.35), emu(1.4), emu(0.35))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'GENIUS BAR'
    run.font.name = 'Calibri'
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = C_PRIMARY

    # Main title
    add_textbox(slide, 0.45, 2.4, 9.0, 1.5,
                'Account Research',
                font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    txb2 = slide.shapes.add_textbox(emu(0.45), emu(3.75), emu(9.0), emu(0.7))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = 'Know your account before you walk in the door.'
    run2.font.name = 'Calibri'
    run2.font.size = Pt(20)
    run2.font.bold = False
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xDD)

    # Bottom label
    add_textbox(slide, 0.45, 6.7, 6.0, 0.35,
                'Genius Bar · April 2026',
                font_size=12, color=RGBColor(0xAA, 0xAA, 0xCC))

    # Red bar at very bottom
    add_rect(slide, 0, 7.38, SLIDE_W_IN, 0.12, fill_color=C_RED)


def build_what_it_is(slide, prs):
    set_slide_bg_white(slide, prs)
    add_header_bars(slide)

    add_textbox(slide, 0.5, 0.28, 12.0, 0.7,
                'Your pre-call research, done in under 5 minutes.',
                font_size=26, bold=True, color=C_DARK_PURPLE)

    body = ('Account Research generates a comprehensive, AI-powered account brief — '
            'grounded in your actual Docusign book of business data — and delivers it '
            'directly to you in Slack. No manual research. No digging through Salesforce. '
            'Just ask, and it arrives.')
    add_textbox(slide, 0.5, 1.15, 9.2, 1.1,
                body, font_size=13, color=C_MUTED_TEXT)

    # Bullet points with purple dots
    bullets = [
        'Walk into any call with a complete picture of your account',
        'Surface expansion whitespace before you\'re in the room',
        'Ground your conversations in data, not assumptions',
    ]
    bullet_top = 2.45
    for bullet in bullets:
        # Purple dot
        dot = add_rect(slide, 0.5, bullet_top + 0.07, 0.12, 0.12, fill_color=C_PRIMARY)
        add_textbox(slide, 0.75, bullet_top, 8.5, 0.4,
                    bullet, font_size=13, color=C_DARK_TEXT)
        bullet_top += 0.48

    # Stat callout box (bottom right)
    stat_box = add_rounded_rect(slide, 10.2, 2.0, 2.7, 1.6,
                                 fill_color=C_LIGHT_PURPLE,
                                 line_color=C_PRIMARY, line_width=1.5, radius_pct=8)
    add_textbox(slide, 10.2, 2.15, 2.7, 0.75,
                '< 5 min', font_size=36, bold=True, color=C_PRIMARY,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, 10.2, 2.9, 2.7, 0.45,
                'from request to link', font_size=11, color=C_MUTED_TEXT,
                align=PP_ALIGN.CENTER)


def build_what_you_get(slide, prs):
    set_slide_bg_white(slide, prs)
    add_header_bars(slide)

    add_textbox(slide, 0.5, 0.28, 12.0, 0.6,
                '9 sections. One document. Everything you need.',
                font_size=26, bold=True, color=C_DARK_PURPLE)

    # 3x3 grid of tiles
    tiles = [
        'Company Profile\n& Financials',
        'Business Performance\n& Strategy',
        'Executive Contacts\n& Tech Stack',
        'Business Map\n& Org Hierarchy',
        'Docusign Footprint',
        'Account Health\nScorecard',
        'Agreement Landscape',
        'Contract Commerce\nEstimates',
        'Priority Map\n& Action Plan',
    ]

    tile_w = 3.8
    tile_h = 1.35
    tile_gap = 0.18
    grid_left = (SLIDE_W_IN - (3 * tile_w + 2 * tile_gap)) / 2
    grid_top = 1.1

    for idx, tile_text in enumerate(tiles):
        row = idx // 3
        col = idx % 3
        tl = grid_left + col * (tile_w + tile_gap)
        tt = grid_top + row * (tile_h + tile_gap)
        shape = add_rounded_rect(slide, tl, tt, tile_w, tile_h,
                                  fill_color=C_LIGHT_PURPLE,
                                  line_color=C_PRIMARY, line_width=1.0, radius_pct=6)
        # Text in tile
        txb = slide.shapes.add_textbox(emu(tl), emu(tt), emu(tile_w), emu(tile_h))
        tf = txb.text_frame
        tf.word_wrap = True
        lines = tile_text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.name = 'Calibri'
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = C_DARK_PURPLE
        # Vertical centering via top margin
        txb.text_frame.margin_top = emu((tile_h - 0.35 * len(lines)) / 2 - 0.05)

    # Callout strip at bottom
    strip_top = grid_top + 3 * tile_h + 2 * tile_gap + 0.18
    strip = add_rounded_rect(slide, 0.5, strip_top, SLIDE_W_IN - 1.0, 0.55,
                              fill_color=C_CREAM_BG,
                              line_color=C_CREAM_BORDER, line_width=1.0, radius_pct=4)
    add_textbox(slide, 0.75, strip_top + 0.1, SLIDE_W_IN - 1.5, 0.4,
                'Account Health and Docusign Footprint are sourced directly from your internal '
                'book of business — not estimated.',
                font_size=11, color=C_AMBER)


def _build_how_to_use_slide(slide, prs, step_num, step_str, heading, body_text, extra_fn=None):
    """Shared structure for slides 4, 5, 6."""
    set_slide_bg_white(slide, prs)
    add_header_bars(slide)

    # Eyebrow
    add_textbox(slide, 0.5, 0.28, 5.0, 0.35,
                'HOW TO USE IT',
                font_size=10, bold=True, color=C_PRIMARY)

    # Decorative large step number
    add_textbox(slide, 0.3, 0.7, 3.0, 2.2,
                step_str,
                font_size=96, bold=True, color=C_LIGHT_PURPLE)

    # Heading
    add_textbox(slide, 0.5, 1.45, 8.5, 0.85,
                heading, font_size=26, bold=True, color=C_DARK_PURPLE)

    # Body
    add_textbox(slide, 0.5, 2.35, 8.5, 1.0,
                body_text, font_size=13, color=C_MUTED_TEXT)

    if extra_fn:
        extra_fn(slide)


def build_step1(slide, prs):
    def extra(slide):
        # Placeholder box for Slack screenshot
        ph = add_rounded_rect(slide, 0.5, 3.55, 12.3, 3.5,
                               fill_color=RGBColor(0xF8, 0xF3, 0xF0),
                               line_color=C_BORDER_LIGHT, line_width=1.5, radius_pct=4)
        add_textbox(slide, 0.5, 3.55, 12.3, 3.5,
                    '[ Slack screenshot ]',
                    font_size=16, color=C_MUTED_TEXT, align=PP_ALIGN.CENTER)

    _build_how_to_use_slide(
        slide, prs,
        step_num=1,
        step_str='01',
        heading='Open the Genius Bar workflow in Slack',
        body_text=('Navigate to the Genius Bar app in Slack and open the Account Research workflow. '
                   'You\'ll find it pinned in the #genius-bar channel.'),
        extra_fn=extra
    )


def build_step2(slide, prs):
    set_slide_bg_white(slide, prs)
    add_header_bars(slide)

    add_textbox(slide, 0.5, 0.28, 5.0, 0.35,
                'HOW TO USE IT', font_size=10, bold=True, color=C_PRIMARY)
    add_textbox(slide, 0.3, 0.7, 3.0, 2.2,
                '02', font_size=96, bold=True, color=C_LIGHT_PURPLE)
    add_textbox(slide, 0.5, 1.45, 10.0, 0.85,
                'Tell it which account you want',
                font_size=26, bold=True, color=C_DARK_PURPLE)

    # Two option boxes side by side
    box_w = 5.7
    box_h = 2.1
    box_top = 2.5

    for col, (label, title, body) in enumerate([
        ('OPTION A', 'Salesforce Account Link',
         'Paste the URL of the account from Salesforce'),
        ('OPTION B', 'GTM_GROUP ID',
         'Enter the GTM_GROUP ID for the account group you want researched'),
    ]):
        bx = 0.5 + col * (box_w + 0.5)
        shape = add_rounded_rect(slide, bx, box_top, box_w, box_h,
                                  fill_color=C_LIGHT_PURPLE,
                                  line_color=C_PRIMARY, line_width=1.5, radius_pct=6)
        add_textbox(slide, bx + 0.2, box_top + 0.18, box_w - 0.4, 0.35,
                    label, font_size=10, bold=True, color=C_PRIMARY)
        add_textbox(slide, bx + 0.2, box_top + 0.52, box_w - 0.4, 0.45,
                    title, font_size=14, bold=True, color=C_DARK_TEXT)
        add_textbox(slide, bx + 0.2, box_top + 1.0, box_w - 0.4, 0.9,
                    body, font_size=12, color=C_MUTED_TEXT)

    add_textbox(slide, 0.5, 4.85, 12.3, 0.55,
                "That's it. One message.",
                font_size=16, italic=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)


def build_step3(slide, prs):
    def extra(slide):
        # Large centered stat
        add_textbox(slide, 1.5, 3.45, 10.3, 0.9,
                    '< 5 minutes',
                    font_size=48, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(slide, 1.5, 4.25, 10.3, 0.45,
                    'from request to link',
                    font_size=14, color=C_MUTED_TEXT, align=PP_ALIGN.CENTER)
        # Placeholder box
        ph = add_rounded_rect(slide, 0.5, 4.85, 12.3, 2.2,
                               fill_color=RGBColor(0xF8, 0xF3, 0xF0),
                               line_color=C_BORDER_LIGHT, line_width=1.5, radius_pct=4)
        add_textbox(slide, 0.5, 4.85, 12.3, 2.2,
                    '[ Slack delivery screenshot ]',
                    font_size=14, color=C_MUTED_TEXT, align=PP_ALIGN.CENTER)

    _build_how_to_use_slide(
        slide, prs,
        step_num=3,
        step_str='03',
        heading='Your brief arrives in Slack in under 5 minutes',
        body_text=('Account Research assembles your brief and delivers a link directly back to you '
                   'in Slack. Open it, share it with your team, use it to prep.'),
        extra_fn=extra
    )


def build_when_to_use(slide, prs):
    set_slide_bg_white(slide, prs)
    add_header_bars(slide)

    add_textbox(slide, 0.5, 0.28, 12.0, 0.6,
                'Use it any time you need to know more about an account.',
                font_size=24, bold=True, color=C_DARK_PURPLE)

    rows_data = [
        (['WHEN', 'WHAT ACCOUNT RESEARCH GIVES YOU'], False, True),
        (['Before a first call',
          'Company profile, org map, exec contacts, tech stack'],
         True, False),
        (['QBR / EBR prep',
          'Account health scorecard, usage trends, renewal proximity'],
         True, False),
        (['Expansion conversation',
          'Whitespace analysis, agreement landscape, product gaps'],
         True, False),
        (['Account planning',
          'Priority map, business unit breakdown, commerce estimates'],
         True, False),
        (['Renewal prep',
          'Health scorecard, contract terms, consumption pacing'],
         True, False),
    ]

    col_widths = [3.5, 9.0]
    add_table_slide(slide, 0.5, 1.1, 12.5, col_widths, rows_data,
                    header_fill=C_DARK_PURPLE, alt_fill=C_ROW_ALT,
                    default_fill=C_WHITE, font_size=13)


def build_quick_reference(slide, prs):
    set_slide_bg_white(slide, prs)
    add_header_bars(slide)

    add_textbox(slide, 0.5, 0.28, 12.0, 0.6,
                'Account Research at a glance',
                font_size=26, bold=True, color=C_DARK_PURPLE)

    col_x = [0.5, 4.8, 9.1]
    col_w = 3.9
    content_top = 0.98

    # Column 1
    add_textbox(slide, col_x[0], content_top, col_w, 0.35,
                'HOW TO REQUEST', font_size=10, bold=True, color=C_PRIMARY)
    steps = [
        '1.  Open Genius Bar in Slack',
        '2.  Paste Salesforce URL or enter GTM_GROUP ID',
        '3.  Receive your link in Slack in < 5 min',
    ]
    for si, step in enumerate(steps):
        add_textbox(slide, col_x[0], content_top + 0.45 + si * 0.52, col_w, 0.5,
                    step, font_size=12, color=C_DARK_TEXT)

    # Divider lines
    for cx in [col_x[1] - 0.15, col_x[2] - 0.15]:
        add_rect(slide, cx, content_top, 0.015, 6.1, fill_color=C_LIGHT_GRAY)

    # Column 2
    add_textbox(slide, col_x[1], content_top, col_w, 0.35,
                'WHAT YOU GET', font_size=10, bold=True, color=C_PRIMARY)
    sections = [
        '1.  Company Profile & Financials',
        '2.  Business Performance & Strategy',
        '3.  Executive Contacts & Tech Stack',
        '4.  Business Map & Org Hierarchy',
        '5.  Docusign Footprint',
        '6.  Account Health Scorecard',
        '7.  Agreement Landscape',
        '8.  Contract Commerce Estimates',
        '9.  Priority Map & Action Plan',
    ]
    for si, sec in enumerate(sections):
        add_textbox(slide, col_x[1], content_top + 0.45 + si * 0.46, col_w, 0.45,
                    sec, font_size=11, color=C_DARK_TEXT)

    # Column 3
    add_textbox(slide, col_x[2], content_top, col_w, 0.35,
                'WHEN TO USE', font_size=10, bold=True, color=C_PRIMARY)
    use_cases = [
        '•  First calls',
        '•  QBRs & EBRs',
        '•  Expansion conversations',
        '•  Account planning',
        '•  Renewal prep',
    ]
    for ui, uc in enumerate(use_cases):
        add_textbox(slide, col_x[2], content_top + 0.45 + ui * 0.46, col_w, 0.45,
                    uc, font_size=12, color=C_DARK_TEXT)

    # Questions label
    q_top = content_top + 0.45 + len(use_cases) * 0.46 + 0.35
    add_textbox(slide, col_x[2], q_top, col_w, 0.35,
                'QUESTIONS?', font_size=10, bold=True, color=C_PRIMARY)
    add_textbox(slide, col_x[2], q_top + 0.35, col_w, 0.4,
                'Reach out in #growth-tools',
                font_size=12, color=C_MUTED_TEXT)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = emu(SLIDE_W_IN)
    prs.slide_height = emu(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]  # blank

    builders = [
        build_cover,
        build_what_it_is,
        build_what_you_get,
        build_step1,
        build_step2,
        build_step3,
        build_when_to_use,
        build_quick_reference,
    ]

    for builder in builders:
        slide = prs.slides.add_slide(blank_layout)
        builder(slide, prs)

    out = '/Users/paul.solans/dev/tools/growth-strategy-2026/context/reference/Account Research - Quick Start.pptx'
    prs.save(out)
    print(f'Saved: {out}')

    import os
    size = os.path.getsize(out)
    print(f'File size: {size:,} bytes ({size/1024:.1f} KB)')


if __name__ == '__main__':
    main()
